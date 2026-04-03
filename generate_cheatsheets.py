#!/usr/bin/env python3
"""Generate PRECC cheatsheet PowerPoint files (EN + ZH)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
import glob
import os

# ---------------------------------------------------------------------------
# Persona definitions
# ---------------------------------------------------------------------------

PERSONAS = [
    {
        "id": "dev",
        "bg": RGBColor(0xFF, 0xF8, 0xF0),
        "accent": RGBColor(0xE8, 0x77, 0x22),
    },
    {
        "id": "tech_vc",
        "bg": RGBColor(0xFF, 0xF9, 0xF0),
        "accent": RGBColor(0xC8, 0x5A, 0x00),
    },
    {
        "id": "biz_vc",
        "bg": RGBColor(0xFF, 0xF5, 0xF0),
        "accent": RGBColor(0xCC, 0x44, 0x00),
    },
    {
        "id": "cto",
        "bg": RGBColor(0xFF, 0xF0, 0xF0),
        "accent": RGBColor(0x99, 0x00, 0x22),
    },
    {
        "id": "devrel",
        "bg": RGBColor(0xFF, 0xF0, 0xEC),
        "accent": RGBColor(0xD9, 0x4F, 0x3D),
    },
]

# ---------------------------------------------------------------------------
# Slide content — English
# ---------------------------------------------------------------------------

SLIDES_EN = [
    # Slide 1 — Developer
    {
        "persona_tag": "For: Developer / Claude Code User",
        "left": {
            "USE CASE": [
                "Claude runs commands from wrong directory → fails → wastes tokens",
                "PRECC intercepts in <3ms and fixes silently",
                "Supports: Rust/cargo, Node/npm, Python/pytest, Go, Make",
            ],
            "DEMO": [
                "1. curl -fsSL https://raw.githubusercontent.com/yijunyu/\n   precc-cc/main/scripts/install.sh | bash",
                "2. precc init",
                "3. precc ingest --all",
                "4. precc report    # see your savings",
            ],
        },
        "right": {
            "RESULTS & BENEFITS": [
                ("34%", "cost reduction on API spend"),
                ("98%", "failures prevented automatically"),
                ("2.93ms", "hook latency — invisible overhead"),
                ("Zero", "config needed — works immediately"),
                ("Learns", "from your own sessions over time"),
            ],
            "INSTALL": [
                "curl -fsSL https://raw.githubusercontent.com/\nperia-ai/precc-cc/main/scripts/install.sh | bash",
            ],
        },
        "tagline": "PRECC — Predictive Error Correction for Claude Code",
    },
    # Slide 2 — Technical VC
    {
        "persona_tag": "For: Technical Investor",
        "left": {
            "USE CASE": [
                "Evaluate PRECC as a technology investment",
                "PreToolUse:Bash hook — canonical SDK integration point",
                "4-pillar architecture: CWD resolution, GDB advisor,\nhistory mining, skill matching",
            ],
            "DEMO": [
                "1. bash demo/demo.sh   # Section 2 — hook timing",
                '2. "2.93ms avg — 0.005% of 60s SDK budget"',
                "3. Section 5 — xxd shows AES-256 encrypted DB header",
                "4. precc report   # real billing metrics",
            ],
        },
        "right": {
            "RESULTS & BENEFITS": [
                ("$296/$878", "saved (34%) — measured, not modeled"),
                ("352/358", "failures prevented (98%)"),
                ("2.93ms", "hook latency — sub-millisecond capable"),
                ("AES-256", "machine-bound encryption (SQLCipher)"),
                ("29", "real sessions measured for baseline"),
            ],
            "MOAT": [
                "Machine-specific heuristics.db — not portable",
                "<5ms hook constraint — tight integration barrier",
                "Fail-open design — PRECC crash never blocks Claude Code",
            ],
        },
        "tagline": "PRECC — Predictive Error Correction for Claude Code",
    },
    # Slide 3 — Non-Technical VC
    {
        "persona_tag": "For: Business Investor",
        "left": {
            "USE CASE": [
                "$878 measured API spend → $296 recoverable waste",
                "34% savings per developer — consistent across sessions",
                "Scales linearly with team size",
                "Zero-friction install — one curl command",
            ],
            "DEMO": [
                "1. bash demo/demo.sh   # Section 1 — red failure lines",
                "2. Section 4 — the numbers",
                '3. "$68/month saved per developer\n   at Claude Code Pro pricing"',
            ],
        },
        "right": {
            "RESULTS & BENEFITS": [
                ("$1,200", "per developer per year in savings"),
                ("$816/yr", "10-person team savings"),
                ("$8,160/yr", "100-person organization savings"),
                ("Day-one", "ROI — no ramp-up period"),
                ("Open-core", "business model — community + premium"),
            ],
            "BUSINESS MODELS": [
                "Per-seat SaaS: $10–20/month",
                "Team heuristics sync (premium)",
                "Skills marketplace",
            ],
        },
        "tagline": "PRECC — Predictive Error Correction for Claude Code",
    },
    # Slide 4 — Enterprise CTO
    {
        "persona_tag": "For: Enterprise CTO / Buyer",
        "left": {
            "USE CASE": [
                "Deploy across dev teams with zero trust footprint",
                "No network calls — fully air-gapped capable",
                "Machine-bound encryption — data stays local",
                "Fail-open design — PRECC crash never blocks Claude Code",
            ],
            "DEMO": [
                "1. Section 0: 'Encryption: AES-256' in init output",
                "2. Section 5: xxd heuristics.db — no SQLite magic bytes",
                "3. precc report — per-engineer ROI for procurement",
                "4. One-command deploy: curl | bash",
            ],
        },
        "right": {
            "RESULTS & BENEFITS": [
                ("98%", "failure prevention rate"),
                ("Zero", "network calls — fully offline"),
                ("AES-256", "SQLCipher encryption"),
                ("Fail-open", "exit 0 on crash — never blocks"),
                ("<5ms", "p99 hook latency"),
            ],
            "SECURITY": [
                "HKDF-SHA256 key from machine-ID + username",
                "No passphrase — no user friction",
                "Reproducible builds — auditable",
                "Linux / macOS / Windows",
            ],
        },
        "tagline": "PRECC — Predictive Error Correction for Claude Code",
    },
    # Slide 5 — Developer Advocate
    {
        "persona_tag": "For: Developer Advocate / Partner",
        "left": {
            "USE CASE": [
                "Wow moment — wrong-dir cargo build fixed in 3ms",
                "vs. 200+ tokens wasted without PRECC",
                "Learning loop: session → mine → skill → prevent",
                "Open skill TOML format for easy integration",
            ],
            "DEMO": [
                "1. Section 2: JSON hook transformation",
                '2. "cargo build → cd /project && cargo build, 2.93ms"',
                "3. Section 3: mining summary — failure-fix pairs",
                "4. precc skills list   # human-readable skill table",
            ],
        },
        "right": {
            "RESULTS & BENEFITS": [
                ("98%", "prevention — compelling demo stat"),
                ("3ms", "wow moment — instant audience hook"),
                ("Smarter", "gets better every session automatically"),
                ("Open TOML", "skill format — hackable & shareable"),
                ("Documented", "heuristics.db schema for partners"),
            ],
            "INTEGRATION": [
                "precc skills export <name>",
                "Custom skill TOML authoring",
                "Ship skill packages for your tool",
            ],
        },
        "tagline": "PRECC — Predictive Error Correction for Claude Code",
    },
]

# ---------------------------------------------------------------------------
# Slide content — Chinese
# ---------------------------------------------------------------------------

SLIDES_ZH = [
    # Slide 1 — Developer
    {
        "persona_tag": "适用人群：开发者 / Claude Code 用户",
        "left": {
            "使用场景": [
                "Claude在错误目录运行命令→失败→浪费Token",
                "PRECC在<3ms内拦截并静默修复",
                "支持：Rust/cargo、Node/npm、Python/pytest、Go、Make",
            ],
            "演示": [
                "1. curl -fsSL https://raw.githubusercontent.com/yijunyu/\n   precc-cc/main/scripts/install.sh | bash",
                "2. precc init",
                "3. precc ingest --all",
                "4. precc report    # 查看节省情况",
            ],
        },
        "right": {
            "效果与收益": [
                ("34%", "API消费成本节省"),
                ("98%", "失败自动预防"),
                ("2.93ms", "钩子延迟——无感知开销"),
                ("零", "配置需求——即装即用"),
                ("自动学习", "从您的会话中持续改进"),
            ],
            "安装": [
                "curl -fsSL https://raw.githubusercontent.com/\nperia-ai/precc-cc/main/scripts/install.sh | bash",
            ],
        },
        "tagline": "PRECC — 面向Claude Code的预测性错误修正",
    },
    # Slide 2 — Technical VC
    {
        "persona_tag": "适用人群：技术投资人",
        "left": {
            "使用场景": [
                "评估PRECC技术投资价值",
                "PreToolUse:Bash钩子——规范SDK集成点",
                "四柱架构：工作目录解析、GDB顾问、\n历史挖掘、技能匹配",
            ],
            "演示": [
                "1. bash demo/demo.sh   # 第2节——钩子计时",
                '2. "2.93ms均值——占60秒SDK预算的0.005%"',
                "3. 第5节——xxd显示AES-256加密DB头部",
                "4. precc report   # 真实计费指标",
            ],
        },
        "right": {
            "效果与收益": [
                ("$296/$878", "节省（34%）——实测数据，非模型推算"),
                ("352/358", "失败预防（98%）"),
                ("2.93ms", "钩子延迟——具备亚毫秒级能力"),
                ("AES-256", "机器绑定加密（SQLCipher）"),
                ("29", "个真实会话用于基准测量"),
            ],
            "护城河": [
                "机器专属heuristics.db——不可移植",
                "<5ms钩子约束——高集成壁垒",
                "故障开放设计——PRECC崩溃不阻断Claude Code",
            ],
        },
        "tagline": "PRECC — 面向Claude Code的预测性错误修正",
    },
    # Slide 3 — Non-Technical VC
    {
        "persona_tag": "适用人群：商业投资人",
        "left": {
            "使用场景": [
                "实测$878 API消费→$296可回收浪费",
                "每位开发者节省34%——跨会话稳定",
                "随团队规模线性扩展",
                "零摩擦安装——一条curl命令",
            ],
            "演示": [
                "1. bash demo/demo.sh   # 第1节——红色失败行",
                "2. 第4节——数字详情",
                '3. "Claude Code Pro定价下\n   每位开发者月均节省$68"',
            ],
        },
        "right": {
            "效果与收益": [
                ("$1,200", "每位开发者每年节省"),
                ("$816/年", "10人团队节省"),
                ("$8,160/年", "百人组织节省"),
                ("当日回报", "无需预热期"),
                ("开源核心", "商业模式——社区版+高级版"),
            ],
            "商业模式": [
                "按席位SaaS：$10–20/月",
                "团队启发式同步（高级功能）",
                "技能市场",
            ],
        },
        "tagline": "PRECC — 面向Claude Code的预测性错误修正",
    },
    # Slide 4 — Enterprise CTO
    {
        "persona_tag": "适用人群：企业CTO / 采购决策者",
        "left": {
            "使用场景": [
                "零信任部署——无网络调用",
                "完全支持离线/隔离网络环境",
                "机器绑定加密——数据始终本地",
                "故障开放设计——PRECC崩溃不阻断Claude Code",
            ],
            "演示": [
                "1. 第0节：初始化输出显示'Encryption: AES-256'",
                "2. 第5节：xxd heuristics.db——无SQLite魔术字节",
                "3. precc report——每工程师ROI采购报告",
                "4. 一键部署：curl | bash",
            ],
        },
        "right": {
            "效果与收益": [
                ("98%", "失败预防率"),
                ("零", "网络调用——完全离线"),
                ("AES-256", "SQLCipher加密"),
                ("故障开放", "崩溃时exit 0——永不阻断"),
                ("<5ms", "P99钩子延迟"),
            ],
            "安全": [
                "HKDF-SHA256密钥来自机器ID+用户名",
                "无口令——零用户摩擦",
                "可复现构建——可审计",
                "Linux / macOS / Windows",
            ],
        },
        "tagline": "PRECC — 面向Claude Code的预测性错误修正",
    },
    # Slide 5 — Developer Advocate
    {
        "persona_tag": "适用人群：开发者倡导者 / 合作伙伴",
        "left": {
            "使用场景": [
                "震撼时刻——错误目录的cargo build在3ms内修复",
                "否则将浪费200+个Token",
                "学习闭环：会话→挖掘→技能→预防",
                "开放技能TOML格式，易于集成",
            ],
            "演示": [
                "1. 第2节：JSON钩子变换",
                '2. "cargo build → cd /project && cargo build，2.93ms"',
                "3. 第3节：挖掘摘要——失败-修复对",
                "4. precc skills list   # 可读技能列表",
            ],
        },
        "right": {
            "效果与收益": [
                ("98%", "预防率——极具说服力的演示数据"),
                ("3ms", "震撼时刻——即时抓住观众"),
                ("越用越聪明", "每个会话自动改进"),
                ("开放TOML", "技能格式——可扩展可分享"),
                ("文档完备", "heuristics.db架构供合作伙伴使用"),
            ],
            "集成方式": [
                "precc skills export <name>",
                "自定义技能TOML编写",
                "为您的工具发布技能包",
            ],
        },
        "tagline": "PRECC — 面向Claude Code的预测性错误修正",
    },
]

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)


def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()

    return shape


def add_text_box(slide, x, y, w, h, text="", font_name="Calibri", font_size=10,
                 bold=False, italic=False, color=None, align=PP_ALIGN.LEFT,
                 word_wrap=True):
    """Add a text box with single paragraph of text."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def add_section_with_bullets(slide, x, y, w, section_title, bullets, accent_color,
                              title_size=11, body_size=10, is_code=False,
                              is_metrics=False):
    """
    Add a section header + bullet list. Returns the bottom y position consumed.
    For is_metrics=True, bullets are (bold_prefix, rest) tuples.
    For is_code=True, bullets are rendered in Courier New.
    """
    LINE_HEIGHT_TITLE = Inches(0.22)
    LINE_HEIGHT_BODY = Inches(0.18)
    SPACING_AFTER_SECTION = Inches(0.12)
    SPACING_BETWEEN_SECTIONS = Inches(0.18)

    cur_y = y

    # Section title
    txb = slide.shapes.add_textbox(x, cur_y, w, LINE_HEIGHT_TITLE)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_title
    run.font.name = "Calibri"
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = accent_color
    cur_y += LINE_HEIGHT_TITLE + Inches(0.04)

    # Bullets
    for bullet in bullets:
        if is_metrics and isinstance(bullet, tuple):
            bold_part, rest_part = bullet
            # Estimate height: allow multi-line
            est_h = LINE_HEIGHT_BODY * 2
            txb = slide.shapes.add_textbox(x, cur_y, w, est_h)
            txb.word_wrap = True
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            # Bold number
            r1 = p.add_run()
            r1.text = bold_part + "  "
            r1.font.name = "Calibri"
            r1.font.size = Pt(body_size)
            r1.font.bold = True
            r1.font.color.rgb = accent_color
            # Normal text
            r2 = p.add_run()
            r2.text = rest_part
            r2.font.name = "Calibri"
            r2.font.size = Pt(body_size)
            r2.font.bold = False
            r2.font.color.rgb = DARK
            cur_y += LINE_HEIGHT_BODY + Inches(0.04)
        else:
            # Plain or code bullet
            text = bullet if isinstance(bullet, str) else str(bullet)
            lines = text.split("\n")
            est_h = LINE_HEIGHT_BODY * max(len(lines), 1) + Inches(0.04)
            txb = slide.shapes.add_textbox(x + Inches(0.15), cur_y, w - Inches(0.15), est_h)
            txb.word_wrap = True
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = ("• " if not is_code else "") + text
            if is_code:
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                run.font.color.rgb = DARK
            else:
                run.font.name = "Calibri"
                run.font.size = Pt(body_size)
                run.font.color.rgb = DARK
            extra_lines = max(len(lines) - 1, 0)
            cur_y += LINE_HEIGHT_BODY + Inches(0.04) + extra_lines * Inches(0.16)

    cur_y += SPACING_AFTER_SECTION
    return cur_y


# ---------------------------------------------------------------------------
# Main slide builder
# ---------------------------------------------------------------------------

def build_slide(prs, persona, slide_data, lang="en", gif_path=None):
    """Build one slide and add it to prs."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg_color = persona["bg"]
    accent = persona["accent"]

    # Background
    set_bg(slide, bg_color)

    # ---- HEADER AREA ----
    HEADER_Y = Inches(0.18)
    LOGO_X = Inches(0.3)
    TAG_X = Inches(8.5)
    TAG_W = Inches(4.5)
    HEADER_H = Inches(0.55)

    # PRECC logo
    add_text_box(slide, LOGO_X, HEADER_Y, Inches(3.0), HEADER_H,
                 text="PRECC", font_name="Calibri", font_size=36,
                 bold=True, color=accent)

    # Persona tag (top-right)
    txb = slide.shapes.add_textbox(TAG_X, HEADER_Y + Inches(0.1), TAG_W, Inches(0.35))
    txb.word_wrap = False
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = slide_data["persona_tag"]
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = accent

    # Horizontal rule
    rule_y = HEADER_Y + HEADER_H + Inches(0.05)
    rule = add_rect(slide,
                    Inches(0.3), rule_y,
                    Inches(12.73), Inches(0.02),
                    fill_color=accent)

    # ---- COLUMNS ----
    COL_Y = rule_y + Inches(0.15)
    LEFT_X = Inches(0.3)
    LEFT_W = Inches(4.8)
    RIGHT_X = Inches(5.5)
    RIGHT_W = Inches(7.3)
    BOTTOM_STRIP_Y = Inches(7.0)  # strip starts here
    MAX_COL_BOTTOM = BOTTOM_STRIP_Y - Inches(0.05)

    # LEFT column
    cur_y_left = COL_Y
    for section_title, items in slide_data["left"].items():
        # Detect if it's a DEMO section (code-style)
        is_code = section_title.upper() in ("DEMO", "演示", "INSTALL", "安装")
        cur_y_left = add_section_with_bullets(
            slide, LEFT_X, cur_y_left, LEFT_W,
            section_title, items, accent,
            title_size=11, body_size=10,
            is_code=is_code,
            is_metrics=False
        )
        if cur_y_left > MAX_COL_BOTTOM:
            break

    # RIGHT column
    cur_y_right = COL_Y
    first_section = True
    for section_title, items in slide_data["right"].items():
        is_metrics = first_section  # first right section is always metrics
        is_code = section_title.upper() in ("INSTALL", "安装")
        cur_y_right = add_section_with_bullets(
            slide, RIGHT_X, cur_y_right, RIGHT_W,
            section_title, items, accent,
            title_size=11, body_size=10,
            is_code=is_code,
            is_metrics=is_metrics
        )
        first_section = False
        if cur_y_right > MAX_COL_BOTTOM:
            break

    # ---- BOTTOM STRIP ----
    STRIP_H = Inches(0.5)
    strip = add_rect(slide,
                     Inches(0), BOTTOM_STRIP_Y,
                     Inches(13.33), STRIP_H,
                     fill_color=accent)

    # Tagline text in strip
    tagline_txb = slide.shapes.add_textbox(
        Inches(0.3), BOTTOM_STRIP_Y + Inches(0.08),
        Inches(9.5), Inches(0.35)
    )
    tagline_txb.word_wrap = False
    tf = tagline_txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_data["tagline"]
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE

    # QR placeholder (white box)
    QR_W = Inches(0.6)
    QR_H = Inches(0.4)
    QR_X = Inches(12.4)
    QR_Y = BOTTOM_STRIP_Y + Inches(0.05)
    qr_box = add_rect(slide, QR_X, QR_Y, QR_W, QR_H, fill_color=WHITE)
    # "QR" label inside
    qr_lbl = slide.shapes.add_textbox(QR_X, QR_Y + Inches(0.1), QR_W, Inches(0.2))
    tf = qr_lbl.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "QR"
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.color.rgb = accent

    # URL next to QR
    url_txb = slide.shapes.add_textbox(
        Inches(10.5), BOTTOM_STRIP_Y + Inches(0.08),
        Inches(1.85), Inches(0.35)
    )
    tf = url_txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "yijunyu.github.io"
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.color.rgb = WHITE

    # ---- ANIMATED GIF (lower-right area) ----
    # Embedded as a video/GIF movie in the right column below content.
    if gif_path and os.path.exists(gif_path):
        GIF_W = Inches(4.5)
        GIF_H = Inches(2.8)
        GIF_X = RIGHT_X + Inches(1.4)
        GIF_Y = BOTTOM_STRIP_Y - GIF_H - Inches(0.1)
        slide.shapes.add_movie(
            gif_path,
            GIF_X, GIF_Y, GIF_W, GIF_H,
            mime_type="image/gif",
        )

    return slide


# ---------------------------------------------------------------------------
# Generate presentations
# ---------------------------------------------------------------------------

def make_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def find_gif(persona_id: str, repo_root: str) -> str | None:
    """Return the most recent GIF for a given persona, or None."""
    pattern = os.path.join(repo_root, f"{persona_id}-*.gif")
    matches = sorted(glob.glob(pattern))
    return matches[-1] if matches else None


def generate(output_path, slides_data, lang="en", repo_root="."):
    prs = make_prs()
    for i, (persona, slide_data) in enumerate(zip(PERSONAS, slides_data)):
        gif_path = find_gif(persona["id"], repo_root)
        if gif_path:
            print(f"  Building slide {i+1}: {slide_data['persona_tag']} + GIF: {os.path.basename(gif_path)}")
        else:
            print(f"  Building slide {i+1}: {slide_data['persona_tag']} (no GIF found)")
        build_slide(prs, persona, slide_data, lang=lang, gif_path=gif_path)
    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    REPO = "/home/y00577373/precc_cc_priv"

    print("Generating English cheatsheet...")
    generate(
        os.path.join(REPO, "precc-cheatsheet-en.pptx"),
        SLIDES_EN,
        lang="en",
        repo_root=REPO,
    )

    print("Generating Chinese cheatsheet...")
    generate(
        os.path.join(REPO, "precc-cheatsheet-zh.pptx"),
        SLIDES_ZH,
        lang="zh",
        repo_root=REPO,
    )

    print("Done.")
